#!/usr/bin/env python3
"""
DawnGuard Hackathon Presentation Generator
Creates a professional PowerPoint presentation for DAWN Black Box hackathon
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Brand colors
DAWN_ORANGE = RGBColor(*hex_to_rgb('#FF6B35'))
DAWN_PURPLE = RGBColor(*hex_to_rgb('#6B35FF'))
DARK_BG = RGBColor(*hex_to_rgb('#1a1a1a'))
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(*hex_to_rgb('#48bb78'))
BLUE = RGBColor(*hex_to_rgb('#4299e1'))
RED = RGBColor(*hex_to_rgb('#e53e3e'))

def create_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background gradient (simulate with shape)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🔐 DawnGuard"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4),
        Inches(8), Inches(1.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Your Family's Private Cloud + AI\nBuilt for DAWN Black Box"
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.color.rgb = DAWN_ORANGE
        paragraph.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(
        Inches(1), Inches(6),
        Inches(8), Inches(1)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = '"Replace Dropbox + ChatGPT. Save $480/year. 100% Private."'
    p = tagline_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def create_problem_slide(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "💸 THE $480/YEAR PROBLEM"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Content box
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    content = """Families Today Pay For:
• Dropbox: $240/year (2TB)
• ChatGPT Plus: $240/year
• Parental Control Software: $50/year

BUT GET:
❌ Zero Privacy (data harvesting)
❌ No Control (terms change anytime)
❌ No AI Safety (kids unsupervised)
❌ Rising Costs (prices always increase)"""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(12)

def create_solution_slide(prs):
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "☀️ THE DAWNGUARD SOLUTION"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "ONE Device. ZERO Monthly Fees. TOTAL Privacy."
    p = tagline_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Features
    features_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    text_frame = features_box.text_frame

    content = """🏠 Family Vault - Unlimited encrypted storage on YOUR hardware

🤖 Private AI Assistant - 100% local, conversations never leave your box

🛡️ AI Guardian - Local content moderation protects kids

🌐 P2P Knowledge Network - Share knowledge, not data

🔗 Blockchain Verified - Solana-powered reputation & auth"""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(16)

def create_innovation_slide(prs):
    """Slide 4: Innovation Highlights"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "💡 WHAT MAKES DAWNGUARD UNIQUE"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Three columns
    col_width = 2.8

    # Column 1: Privacy
    col1_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(col_width), Inches(4.5))
    text_frame = col1_box.text_frame
    content1 = """🔐 PRIVACY
• Zero-Knowledge Proof Auth
  First AI using ZKP

• Multi-Layer Encryption
  5 security layers

• Local AI Moderation
  NO competitor does this"""
    text_frame.text = content1
    for p in text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    # Column 2: AI
    col2_box = slide.shapes.add_textbox(Inches(3.6), Inches(1.8), Inches(col_width), Inches(4.5))
    text_frame = col2_box.text_frame
    content2 = """🤖 AI INNOVATION
• Graceful Degradation
  Works without Ollama

• Streaming Responses
  Real-time chat

• AI-Powered Features
  Smart search, tagging"""
    text_frame.text = content2
    for p in text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    # Column 3: Family
    col3_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.8), Inches(col_width), Inches(4.5))
    text_frame = col3_box.text_frame
    content3 = """👨‍👩‍👧‍👦 FAMILY FIRST
• Only for Families
  Not enterprise adapted

• Parental Controls
  Activity monitoring

• Storage Quotas
  Digital responsibility"""
    text_frame.text = content3
    for p in text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

def create_cypherpunk_slide(prs):
    """Slide 5: Cypherpunk Values"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "☀️ PRAISE THE SUN: CYPHERPUNK AUTHENTICITY"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    quote_frame = quote_box.text_frame
    quote_frame.text = '"Privacy is necessary for an open society in the electronic age."'
    p = quote_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3.5))
    text_frame = content_box.text_frame
    content = """✅ Privacy by Design - Real cryptography: AES-256, RSA-2048, ZKP

✅ User Sovereignty - Your data, your hardware, your keys

✅ Decentralization - P2P mesh, no central server

✅ Open Source - Audit the code yourself

✅ Economic Incentives - Reputation rewards contribution

"Cypherpunks write code." - We built it."""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(14)

def create_blackbox_fit_slide(prs):
    """Slide 6: Perfect Fit for Black Box"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "⚡ WHY THIS BELONGS ON EVERY BLACK BOX"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    text_frame = content_box.text_frame
    content = """🏠 ALWAYS-ON HOME SERVER
✓ AI runs 24/7 for family  ✓ P2P node always available

💾 LOCAL STORAGE
✓ Unlimited file storage (add drives)  ✓ Zero monthly costs

⚡ LOCAL COMPUTE
✓ Run LLMs locally (Llama 3.2)  ✓ Content moderation on-device

🌐 NETWORK INTEGRATION
✓ Meshes with DAWN network  ✓ P2P to other Black Boxes

💰 ROI PROOF
✓ Saves $480/year  ✓ Black Box pays for itself in < 1 year"""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(18)

def create_architecture_slide(prs):
    """Slide 7: Technical Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🏗️ TECHNICAL ARCHITECTURE"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Architecture diagram (text-based)
    arch_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(3))
    text_frame = arch_box.text_frame
    arch_text = """┌─ YOUR BLACK BOX ─────────────┐
│  Django Web App ←→ Ollama AI  │
│         ↕                      │
│  Encryption Layer (AES-256)   │
│         ↕                      │
│  SQLite Database (Encrypted)  │
└───────────────────────────────┘
         ↕              ↕
   Solana Devnet    P2P Mesh"""

    text_frame.text = arch_text
    for p in text_frame.paragraphs:
        p.font.name = 'Courier New'
        p.font.size = Pt(18)
        p.font.color.rgb = GREEN

    # Stack info
    stack_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.5))
    text_frame = stack_box.text_frame
    stack_text = """Stack: Django 5.0, Python 3.11, Ollama (Llama 3.2), Solana Devnet
Encryption: AES-256, RSA-2048, PBKDF2 | Deploy: Docker Compose"""
    text_frame.text = stack_text
    for p in text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

def create_impact_slide(prs):
    """Slide 8: Impact & Market"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🌍 REAL-WORLD IMPACT"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    text_frame = content_box.text_frame
    content = """Target Market:
👨‍👩‍👧‍👦 4.3M families in US pay for Dropbox | 💰 $15B cloud storage market

DawnGuard Impact per Family:
💵 $480 saved/year | 🌱 40kg CO₂ saved | 🔐 100% privacy

Competitive Advantages:
✅ ONLY family-focused privacy app
✅ ONLY local AI content moderation
✅ ONLY ZKP authentication for AI
✅ ONLY unlimited storage for $0

Network Effects: Each Black Box strengthens P2P mesh"""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(16)

def create_production_slide(prs):
    """Slide 9: Production Ready"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 PRODUCTION-READY CODE"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Deployment
    deploy_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.5))
    text_frame = deploy_box.text_frame
    deploy_text = """One-Command Deploy:
$ git clone https://github.com/shariqazeem/DawnGuard.git
$ cd DawnGuard && ./scripts/setup.sh"""
    text_frame.text = deploy_text
    for p in text_frame.paragraphs:
        p.font.name = 'Courier New'
        p.font.size = Pt(18)
        p.font.color.rgb = GREEN

    # Stats
    stats_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(3))
    text_frame = stats_box.text_frame
    content = """Code Statistics:
✅ 25+ Django models | 30+ view functions | 120KB code
✅ 32 HTML templates | 5 crypto handlers | 13 migrations

Quality:
• Clean MVC architecture | Comprehensive error handling
• CSRF protection | SQL injection prevention
• Responsive UI (mobile-ready) | Docker containerized"""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(14)

def create_roadmap_slide(prs):
    """Slide 10: Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🛣️ ROADMAP"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DAWN_ORANGE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5))
    text_frame = content_box.text_frame
    content = """PHASE 1: MVP ✅ (Current - Hackathon)
• All 9 core features implemented and working
• Family Vault, AI Assistant, AI Guardian, P2P Network, Blockchain

PHASE 2: ENHANCEMENT (Q1 2026)
• NFT badges to mainnet | DAWN token integration
• Mobile app | Voice assistant | Advanced AI models

PHASE 3: ECOSYSTEM (Q2-Q3 2026)
• DAO governance | Federated AI training | AI marketplace
• Plugin system | Enterprise features

PHASE 4: SCALE (Q4 2026+)
• 100K+ Black Boxes | Token economy | Cross-box AI inference"""

    text_frame.text = content
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(12)

def create_closing_slide(prs):
    """Slide 11: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "☀️ PRAISE THE SUN\nDawnGuard: Privacy for Every Family"
    for p in title_frame.paragraphs:
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = DAWN_ORANGE
        p.alignment = PP_ALIGN.CENTER

    # Summary
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    text_frame = summary_box.text_frame
    summary = """✅ Innovation: ZKP + Local AI + P2P + Blockchain
✅ Technical: Production-ready, 120KB code, real crypto
✅ Impact: $480 saved/year, 100% privacy, family safety
✅ Clarity: Replace Dropbox + ChatGPT with YOUR Black Box"""

    text_frame.text = summary
    for p in text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    # Why we win
    win_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(1.5))
    text_frame = win_box.text_frame
    win_text = """Why DawnGuard Wins:
🏆 Only family-focused privacy app | 🏆 Only local AI moderation
🏆 Only ZKP auth for AI | 🏆 Perfect Black Box fit

One Family, One Box, Zero Compromises."""

    text_frame.text = win_text
    for p in text_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = GREEN
        p.font.bold = True

def main():
    """Create the presentation"""
    print("Creating DawnGuard Hackathon Presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create all slides
    create_title_slide(prs)
    print("✓ Slide 1: Title")

    create_problem_slide(prs)
    print("✓ Slide 2: Problem")

    create_solution_slide(prs)
    print("✓ Slide 3: Solution")

    create_innovation_slide(prs)
    print("✓ Slide 4: Innovation")

    create_cypherpunk_slide(prs)
    print("✓ Slide 5: Cypherpunk Values")

    create_blackbox_fit_slide(prs)
    print("✓ Slide 6: Black Box Fit")

    create_architecture_slide(prs)
    print("✓ Slide 7: Architecture")

    create_impact_slide(prs)
    print("✓ Slide 8: Impact & Market")

    create_production_slide(prs)
    print("✓ Slide 9: Production Ready")

    create_roadmap_slide(prs)
    print("✓ Slide 10: Roadmap")

    create_closing_slide(prs)
    print("✓ Slide 11: Call to Action")

    # Save
    output_file = 'DawnGuard_Hackathon_Presentation.pptx'
    prs.save(output_file)
    print(f"\n✅ Presentation saved: {output_file}")
    print(f"📊 Total slides: 11")
    print(f"⏱️  Presentation time: ~8-10 minutes")
    print("\n🎯 Ready for your demo! Good luck! 🚀")

if __name__ == '__main__':
    main()
